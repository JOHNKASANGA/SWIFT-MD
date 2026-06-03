import io
import os
import httpx
from pypdf import PdfReader
import re
from bs4 import BeautifulSoup


def download_from_url(file_url: str) -> bytes:
    """Download file content, handling Google Drive and MediaFire special cases."""
    
    # Google Drive detection
    if "drive.google.com" in file_url or "docs.google.com" in file_url:
        return download_google_drive(file_url)
    
    # MediaFire detection
    if "mediafire.com" in file_url:
        return download_mediafire(file_url)
    
    # Normal direct download
    response = httpx.get(file_url, timeout=60.0, follow_redirects=True)
    response.raise_for_status()
    return response.content


def download_google_drive(url: str) -> bytes:
    """Handle Google Drive downloads including virus scan confirmation."""
    
    # Extract file ID from various Google Drive URL formats
    file_id = None
    patterns = [
        r'/file/d/([a-zA-Z0-9_-]+)',
        r'id=([a-zA-Z0-9_-]+)',
        r'/d/([a-zA-Z0-9_-]+)',
    ]
    for pattern in patterns:
        match = re.search(pattern, url)
        if match:
            file_id = match.group(1)
            break
    
    if not file_id:
        raise Exception(f"Could not extract Google Drive file ID from URL: {url}")
    
    # Try direct download first
    direct_url = f"https://drive.google.com/uc?export=download&id={file_id}"
    response = httpx.get(direct_url, timeout=60.0, follow_redirects=True)
    
    content_type = response.headers.get("content-type", "")
    
    # If we got HTML, it's the virus scan confirmation page
    if "text/html" in content_type:
        soup = BeautifulSoup(response.text, "html.parser")
        form = soup.find("form")
        if form:
            action = form.get("action")
            inputs = {}
            for inp in form.find_all("input"):
                name = inp.get("name")
                if name:
                    inputs[name] = inp.get("value", "")
            
            # Submit the confirmation form
            confirm_response = httpx.post(
                action,
                data=inputs,
                timeout=120.0,
                follow_redirects=True
            )
            confirm_response.raise_for_status()
            return confirm_response.content
    
    return response.content


def download_mediafire(url: str) -> bytes:
    """Handle MediaFire downloads by extracting direct download link."""
    
    response = httpx.get(url, timeout=30.0, follow_redirects=True)
    soup = BeautifulSoup(response.text, "html.parser")
    
    download_btn = soup.find("a", {"id": "downloadButton"})
    if download_btn:
        direct_url = download_btn.get("href")
        file_response = httpx.get(direct_url, timeout=120.0, follow_redirects=True)
        file_response.raise_for_status()
        return file_response.content
    
    raise Exception("Could not find MediaFire download link")

def extract_text_from_url(file_url: str) -> str:
    """Download a file from URL and extract text based on format."""
    
    content = download_from_url(file_url)
    
    lower_url = file_url.lower()
    
    if lower_url.endswith(".docx"):
        return extract_docx(content)
    elif lower_url.endswith(".pptx"):
        return extract_pptx(content)
    else:
        return extract_pdf(content)

def extract_pdf(content: bytes) -> str:
    pdf_stream = io.BytesIO(content)
    reader = PdfReader(pdf_stream)
    text_parts = []
    for page in reader.pages:
        text = page.extract_text()
        if text:
            text_parts.append(text)
    return "\n\n".join(text_parts)


def extract_docx(content: bytes) -> str:
    from docx import Document
    doc_stream = io.BytesIO(content)
    doc = Document(doc_stream)
    text_parts = []
    for para in doc.paragraphs:
        if para.text.strip():
            text_parts.append(para.text)
    return "\n\n".join(text_parts)


def extract_pptx(content: bytes) -> str:
    from pptx import Presentation
    pptx_stream = io.BytesIO(content)
    prs = Presentation(pptx_stream)
    text_parts = []
    for slide in prs.slides:
        slide_text = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    if para.text.strip():
                        slide_text.append(para.text)
        if slide_text:
            text_parts.append("\n".join(slide_text))
    return "\n\n".join(text_parts)